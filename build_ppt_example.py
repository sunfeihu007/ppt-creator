#!/usr/bin/env python3
"""
PPT组装示例脚本
将生成的图片组装成PPTX文件

使用方法:
1. 确保已安装 python-pptx: pip install python-pptx
2. 运行脚本: python build_ppt_example.py
"""

import os
from pptx import Presentation
from pptx.util import Inches

IMAGE_DIR = "./ppt_pages"
OUTPUT_PPTX = "./output.pptx"

def build_ppt():
    """组装PPT"""
    print("=" * 60)
    print("PPT组装")
    print("=" * 60)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9比例
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    image_files = sorted([f for f in os.listdir(IMAGE_DIR) if f.endswith('.png')])
    print(f"找到 {len(image_files)} 张图片")
    
    for i, image_file in enumerate(image_files, start=1):
        image_path = os.path.join(IMAGE_DIR, image_file)
        
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            image_path,
            0, 0,
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        print(f"  ✓ 添加页面 {i}: {image_file}")
    
    prs.save(OUTPUT_PPTX)
    print(f"\nPPT已保存: {OUTPUT_PPTX}")
    print(f"总页数: {len(image_files)}")

if __name__ == "__main__":
    build_ppt()
