#!/usr/bin/env python3
"""PPTX 组装器：gate检查 → 图片压缩 → 按plan顺序组装 → 注入演讲者备注。

用法: build_ppt.py [--allow-generated] [--no-compress] [--out 路径]
gate: 默认要求全部页面 status=approved（--allow-generated 放宽到 generated，仅调试用）。
压缩: >2MB 的 PNG 转 JPEG(q85)，可缩小5-20倍体积，便于分发。
"""
import argparse
import io
import json
import os
import re
import sys

WS = os.environ.get("PPTC_WORKSPACE", "./ppt_workspace")
STATUS_ORDER = ["pending", "prompted", "generated", "approved"]


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--allow-generated", action="store_true")
    ap.add_argument("--no-compress", action="store_true")
    ap.add_argument("--out")
    args = ap.parse_args()

    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    plan = json.load(open(os.path.join(WS, "plan.json"), encoding="utf-8"))

    # ---- GATE：拒绝组装未完成的PPT（防跳步的最后一道闸）----
    need = "generated" if args.allow_generated else "approved"
    lvl = STATUS_ORDER.index(need)
    bad = [f"{p['id']}({p['status']})" for p in plan["pages"]
           if STATUS_ORDER.index(p["status"]) < lvl]
    if bad:
        sys.exit(f"[build] GATE FAILED —— 以下页面未达 {need}，按七步流程先完成它们:\n  "
                 + ", ".join(bad))
    missing_notes = [p["id"] for p in plan["pages"] if not p.get("notes", "").strip()]
    if missing_notes:
        print(f"[build] 警告：以下页面缺演讲者备注: {missing_notes}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    total_in = total_out = 0
    for p in plan["pages"]:
        path = os.path.join(WS, p["image"])
        size_in = os.path.getsize(path)
        total_in += size_in
        stream = path
        if not args.no_compress and size_in > 2 * 1024 * 1024:
            im = Image.open(path).convert("RGB")
            buf = io.BytesIO()
            im.save(buf, "JPEG", quality=85)
            buf.seek(0)
            stream = buf
            total_out += buf.getbuffer().nbytes
        else:
            total_out += size_in
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(stream, 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        if p.get("notes"):
            slide.notes_slide.notes_text_frame.text = p["notes"]

    topic = re.sub(r'[\\/:*?"<>|]', "_", plan.get("topic") or "presentation")
    out = args.out or os.path.join(WS, "output", f"{topic}.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    mb = os.path.getsize(out) / 1048576
    print(f"[build] ✓ {out}")
    print(f"[build] {len(plan['pages'])} 页 | 图片 {total_in/1048576:.0f}MB -> "
          f"约{total_out/1048576:.0f}MB | 成品 {mb:.1f}MB")


if __name__ == "__main__":
    main()
